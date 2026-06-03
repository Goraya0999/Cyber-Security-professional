from __future__ import annotations

import os
import zipfile
from pathlib import Path
from xml.sax.saxutils import escape
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "Distributed_Intrusion_Detection_System_Defense_Muhammad_Zaid.pptx"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(6, 182, 212)
GREEN = RGBColor(16, 185, 129)
RED = RGBColor(239, 68, 68)
AMBER = RGBColor(245, 158, 11)
BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(203, 213, 225)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
MONO_FONT = "Cascadia Mono"


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(slide, x, y, w, h, text="", font_size=16, color=SLATE, bold=False,
                align=PP_ALIGN.LEFT, font=BODY_FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.65, 0.35, 8.3, 0.45, title, 25, NAVY, True, font=TITLE_FONT)
    if subtitle:
        add_textbox(slide, 0.68, 0.85, 8.3, 0.3, subtitle, 9.5, MUTED)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.2), Inches(1.15), Inches(0.04))
    set_fill(bar, BLUE)
    bar.line.fill.background()


def add_footer(slide, num):
    add_textbox(slide, 0.65, 7.12, 5.5, 0.18, "Distributed Intrusion Detection System - University Project Defense", 7.5, MUTED)
    add_textbox(slide, 12.0, 7.1, 0.45, 0.2, f"{num:02d}", 8, MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, font_size=13, color=SLATE, gap=0.08):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(gap * 40)
        p._p.get_or_add_pPr().set("marL", "260000")
        p._p.get_or_add_pPr().set("hanging", "120000")
    return tf


def card(slide, x, y, w, h, title=None, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(s, fill)
    set_line(s, line, 1)
    if title:
        add_textbox(slide, x + 0.18, y + 0.14, w - 0.36, 0.25, title, 10, NAVY, True)
    return s


def pill(slide, x, y, w, text, color, bg=None):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    set_fill(s, bg or RGBColor(239, 246, 255))
    set_line(s, color, 0.8)
    tf = s.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = BODY_FONT
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = color
    return s


def icon_circle(slide, x, y, label, color=BLUE, text_color=WHITE, size=0.48):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    set_fill(s, color)
    s.line.fill.background()
    tf = s.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = TITLE_FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = text_color
    return s


def arrow(slide, x1, y1, x2, y2, color=BLUE, width=2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def add_metric(slide, x, y, label, value, accent):
    card(slide, x, y, 2.45, 0.86)
    icon_circle(slide, x + 0.16, y + 0.18, "", accent, size=0.18)
    add_textbox(slide, x + 0.42, y + 0.16, 1.8, 0.2, label.upper(), 7.5, MUTED, True)
    add_textbox(slide, x + 0.42, y + 0.42, 1.8, 0.3, value, 18, NAVY, True, font=TITLE_FONT)


def add_code(slide, x, y, w, h, code):
    s = card(slide, x, y, w, h, fill=RGBColor(15, 23, 42), line=RGBColor(30, 41, 59), radius=True)
    add_textbox(slide, x + 0.2, y + 0.12, w - 0.4, 0.2, "scan.ts / detection_service.py", 8, RGBColor(148, 163, 184), True)
    box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.42), Inches(w - 0.35), Inches(h - 0.55))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    for idx, line in enumerate(code.splitlines()):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = MONO_FONT
        r.font.size = Pt(8.2)
        r.font.color.rgb = RGBColor(226, 232, 240)
    return s


def add_slide_notes(pptx_path: Path, notes: list[str]):
    # python-pptx does not expose notes slides, so add simple OOXML note parts.
    ns_ct = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}
    ns_rel = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
    replacements: dict[str, bytes] = {}
    note_parts: dict[str, bytes] = {}

    with zipfile.ZipFile(pptx_path, "r") as z:
        ct_root = ET.fromstring(z.read("[Content_Types].xml"))
        existing = {el.attrib.get("PartName") for el in ct_root}
        for i, note in enumerate(notes, start=1):
            part = f"/ppt/notesSlides/notesSlide{i}.xml"
            if part not in existing:
                ET.SubElement(
                    ct_root,
                    f"{{{ns_ct['ct']}}}Override",
                    {
                        "PartName": part,
                        "ContentType": "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml",
                    },
                )

            note_lines = "".join(
                f"<a:p><a:r><a:rPr lang=\"en-US\" sz=\"1200\"/><a:t>{escape(line)}</a:t></a:r></a:p>"
                for line in note.splitlines()
            )
            notes_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Notes Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr/>
        <p:txBody><a:bodyPr/><a:lstStyle/>{note_lines}</p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>"""
            note_parts[f"ppt/notesSlides/notesSlide{i}.xml"] = notes_xml.encode("utf-8")

            rel_path = f"ppt/slides/_rels/slide{i}.xml.rels"
            rel_root = ET.fromstring(z.read(rel_path))
            if not any(el.attrib.get("Type", "").endswith("/notesSlide") for el in rel_root):
                ids = [el.attrib.get("Id", "") for el in rel_root]
                max_id = 0
                for rid in ids:
                    if rid.startswith("rId") and rid[3:].isdigit():
                        max_id = max(max_id, int(rid[3:]))
                ET.SubElement(
                    rel_root,
                    f"{{{ns_rel['rel']}}}Relationship",
                    {
                        "Id": f"rId{max_id + 1}",
                        "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
                        "Target": f"../notesSlides/notesSlide{i}.xml",
                    },
                )
            replacements[rel_path] = ET.tostring(rel_root, encoding="utf-8", xml_declaration=True)

        replacements["[Content_Types].xml"] = ET.tostring(ct_root, encoding="utf-8", xml_declaration=True)
        replacements.update(note_parts)

        tmp_path = pptx_path.with_suffix(".tmp.pptx")
        with zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as out:
            skip = set(replacements)
            seen = set()
            for info in z.infolist():
                if info.filename in skip or info.filename in seen:
                    continue
                seen.add(info.filename)
                out.writestr(info, z.read(info.filename))
            for name, data in replacements.items():
                out.writestr(name, data)
    tmp_path.replace(pptx_path)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    notes: list[str] = []

    def new_slide(title=None, subtitle=None):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = BG
        if title:
            add_title(slide, title, subtitle)
            add_footer(slide, len(prs.slides))
        return slide

    # 1 Title
    slide = new_slide()
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(4.35), Inches(7.5))
    set_fill(left, NAVY)
    left.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4.35), Inches(0), Inches(0.12), Inches(7.5))
    set_fill(accent, BLUE)
    accent.line.fill.background()
    icon_circle(slide, 0.75, 0.65, "D", BLUE, size=0.68)
    add_textbox(slide, 0.75, 1.55, 3.0, 0.28, "Academic Project Defense", 12, RGBColor(191, 219, 254), True)
    add_textbox(slide, 4.9, 1.15, 7.1, 1.45, "Distributed Intrusion\nDetection System", 35, NAVY, True, font=TITLE_FONT)
    add_textbox(slide, 4.95, 2.82, 6.7, 0.4, "AI-assisted monitoring, detection, alerts, and analytics for network traffic", 15, SLATE)
    pill(slide, 4.95, 3.45, 1.35, "React", BLUE)
    pill(slide, 6.45, 3.45, 1.35, "FastAPI", GREEN, RGBColor(236, 253, 245))
    pill(slide, 7.95, 3.45, 1.65, "Random Forest", RED, RGBColor(254, 242, 242))
    pill(slide, 9.8, 3.45, 1.65, "PostgreSQL", CYAN, RGBColor(236, 254, 255))
    card(slide, 4.95, 4.35, 6.8, 1.2)
    add_textbox(slide, 5.25, 4.58, 6.0, 0.28, "Student: Muhammad Zaid", 14, NAVY, True)
    add_textbox(slide, 5.25, 4.95, 6.0, 0.25, "Roll No: 232674  |  Member: M. Shafiq (232699)", 12, SLATE)
    add_textbox(slide, 5.25, 5.28, 6.0, 0.25, "GCUF", 12, SLATE, True)
    notes.append("Introduce the project title and team. Emphasize that the system combines a dashboard, API services, persistent logging, alerts, and an ML detection model.")

    # 2 Overview
    slide = new_slide("Project Overview", "Introduction, problem statement, and objectives")
    card(slide, 0.75, 1.55, 3.75, 4.55, "Introduction")
    add_bullets(slide, [
        "A distributed intrusion detection system monitors traffic across multiple nodes.",
        "Traffic is classified as benign or malicious using feature-based detection.",
        "Security events are stored, visualized, and converted into actionable alerts.",
    ], 0.95, 2.05, 3.25, 2.3)
    card(slide, 4.85, 1.55, 3.75, 4.55, "Problem Statement")
    add_bullets(slide, [
        "Manual log review is slow and error-prone.",
        "Single-node IDS designs miss distributed or high-volume attack behavior.",
        "Teams need a real-time view of threats, nodes, and incident history.",
    ], 5.05, 2.05, 3.25, 2.3)
    card(slide, 8.95, 1.55, 3.75, 4.55, "Objectives")
    add_bullets(slide, [
        "Detect DoS, DDoS, PortScan, and brute-force attack classes.",
        "Provide dashboard analytics, recent logs, and alert resolution workflow.",
        "Support authentication, API contracts, testing, and containerized deployment.",
    ], 9.15, 2.05, 3.25, 2.5)
    notes.append("Explain the project as a practical defense dashboard. Keep the focus on why distributed monitoring and rapid classification are needed.")

    # 3 Existing System
    slide = new_slide("Existing System & Challenges", "Limitations that motivated the proposed DIDS")
    add_textbox(slide, 0.8, 1.48, 4.8, 0.28, "Traditional / current approach", 15, NAVY, True)
    rows = [
        ("Fragmented logs", "Logs remain scattered across machines and services."),
        ("Delayed response", "Analysts discover attacks after damage or downtime."),
        ("Signature dependence", "Simple rules miss variations of known behavior."),
        ("Weak visibility", "No unified picture of node health, threat rate, and trends."),
    ]
    for i, (t, d) in enumerate(rows):
        y = 2.0 + i * 0.9
        icon_circle(slide, 0.95, y, str(i + 1), RED if i < 2 else AMBER, size=0.36)
        add_textbox(slide, 1.45, y - 0.02, 3.0, 0.24, t, 12, NAVY, True)
        add_textbox(slide, 1.45, y + 0.25, 4.3, 0.22, d, 9.5, SLATE)
    card(slide, 6.5, 1.55, 5.8, 4.55, "Need for Proposed Solution")
    add_bullets(slide, [
        "Centralized event collection from multiple network nodes.",
        "Automated detection with confidence score and threat category.",
        "Real-time dashboard for operational awareness.",
        "Persistent audit trail for reports, investigation, and viva demonstration.",
        "Authentication and controlled API access for safer usage.",
    ], 6.8, 2.05, 5.15, 3.15)
    notes.append("Contrast the old workflow with the proposed workflow. Highlight speed, visibility, consistency, and the ability to preserve evidence.")

    # 4 Proposed Solution
    slide = new_slide("Proposed Solution", "A full-stack DIDS with detection, alerts, analytics, and history")
    center = card(slide, 4.75, 2.05, 3.8, 1.1, "DIDS Platform", fill=RGBColor(239, 246, 255), line=BLUE)
    add_textbox(slide, 5.05, 2.48, 3.1, 0.22, "Dashboard + API + ML Detection", 12, BLUE, True, align=PP_ALIGN.CENTER)
    features = [
        ("Packet Ingestion", "Submit manual traffic or tcpdump-style packet metadata from sensor nodes."),
        ("Threat Alerts", "Automatically creates alerts for malicious detections."),
        ("Analytics", "Traffic trends, threat distribution, and dashboard summary."),
        ("Audit Logs", "Paginated logs with filters, statuses, and confidence score."),
        ("Auth", "JWT user authentication and protected routes."),
        ("Deployment", "Docker Compose with frontend, backend, and PostgreSQL."),
    ]
    positions = [(0.8, 1.65), (0.8, 4.2), (4.75, 4.6), (8.95, 1.65), (8.95, 4.2), (4.75, 0.95)]
    colors = [BLUE, RED, CYAN, GREEN, AMBER, SLATE]
    for idx, ((title, desc), (x, y)) in enumerate(zip(features, positions)):
        card(slide, x, y, 3.2, 1.15)
        icon_circle(slide, x + 0.18, y + 0.22, title[0], colors[idx], size=0.42)
        add_textbox(slide, x + 0.72, y + 0.18, 2.25, 0.23, title, 11.5, NAVY, True)
        add_textbox(slide, x + 0.72, y + 0.48, 2.2, 0.42, desc, 8.3, SLATE)
        arrow(slide, x + (3.2 if x < 4.75 else 0), y + 0.57, 4.75 if x < 4.75 else 8.55, 2.6, colors[idx], 1.4)
    notes.append("Describe the solution as an integrated platform. Mention that the key feature is not only prediction but the complete response workflow around it.")

    # 5 Architecture
    slide = new_slide("System Architecture", "Main components and data movement")
    layers = [
        (0.85, 1.6, 2.3, 1.0, "Users / Admin", "Browser-based security operator", BLUE),
        (3.65, 1.6, 2.3, 1.0, "React Dashboard", "Vite, React Query, charts", CYAN),
        (6.45, 1.6, 2.3, 1.0, "API Gateway", "Express routes and contracts", GREEN),
        (9.25, 1.6, 2.3, 1.0, "PostgreSQL", "Logs, alerts, users", SLATE),
    ]
    for x, y, w, h, title, desc, color in layers:
        card(slide, x, y, w, h, fill=WHITE, line=color)
        icon_circle(slide, x + 0.15, y + 0.23, title[0], color, size=0.42)
        add_textbox(slide, x + 0.68, y + 0.2, 1.45, 0.24, title, 10.5, NAVY, True)
        add_textbox(slide, x + 0.68, y + 0.5, 1.45, 0.27, desc, 7.8, MUTED)
    arrow(slide, 3.15, 2.1, 3.65, 2.1)
    arrow(slide, 5.95, 2.1, 6.45, 2.1)
    arrow(slide, 8.75, 2.1, 9.25, 2.1)
    card(slide, 3.0, 3.6, 2.45, 1.2, "Network Sensors", fill=RGBColor(240, 253, 250), line=GREEN)
    add_textbox(slide, 3.18, 4.08, 2.1, 0.22, "packet metadata + tcpdump input", 8.5, GREEN, True, align=PP_ALIGN.CENTER)
    card(slide, 6.05, 3.55, 3.2, 1.35, "ML Detection Service", fill=RGBColor(254, 242, 242), line=RED)
    add_textbox(slide, 6.35, 4.03, 2.6, 0.23, "FastAPI + Random Forest model", 10.3, RED, True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 6.35, 4.35, 2.6, 0.2, "79 CICIDS-style flow features", 8.2, MUTED, align=PP_ALIGN.CENTER)
    arrow(slide, 4.25, 3.6, 6.7, 2.6, GREEN)
    arrow(slide, 7.65, 3.55, 7.65, 2.6, RED)
    arrow(slide, 8.0, 3.55, 9.65, 2.6, RED)
    card(slide, 0.85, 5.45, 10.7, 0.65, "Component Explanation")
    add_bullets(slide, [
        "Frontend collects scan inputs and visualizes health, traffic, logs, and alerts.",
        "API layer validates requests, persists records, and returns typed responses.",
        "ML service classifies traffic and maps malicious classes to severity and alert metadata.",
    ], 1.05, 5.72, 10.1, 0.48, font_size=8.6)
    notes.append("Walk through data movement from a user action to API validation, model inference, database persistence, and dashboard visualization.")

    # 6 Technologies
    slide = new_slide("Technologies Used", "Tools selected for a maintainable academic prototype")
    techs = [
        ("Frontend", "React, Vite, TypeScript, Tailwind CSS, Recharts, Lucide icons"),
        ("Backend API", "Node.js, Express, OpenAPI contract, Zod validation"),
        ("ML API", "Python, FastAPI, scikit-learn Random Forest, NumPy, joblib"),
        ("Database", "PostgreSQL, SQLAlchemy models, Drizzle schema"),
        ("Security", "JWT authentication, password hashing, protected routes/API keys"),
        ("DevOps & Tests", "Docker Compose, Vitest, Pytest, generated API clients"),
    ]
    for i, (title, desc) in enumerate(techs):
        x = 0.85 + (i % 2) * 6.0
        y = 1.45 + (i // 2) * 1.45
        card(slide, x, y, 5.35, 1.05)
        icon_circle(slide, x + 0.18, y + 0.24, title[0], [BLUE, GREEN, RED, CYAN, AMBER, SLATE][i], size=0.42)
        add_textbox(slide, x + 0.75, y + 0.18, 2.1, 0.24, title, 11.5, NAVY, True)
        add_textbox(slide, x + 0.75, y + 0.5, 4.25, 0.3, desc, 8.7, SLATE)
    notes.append("Mention why the stack is appropriate: React for UI, FastAPI/scikit-learn for ML, PostgreSQL for durable records, and Docker for reproducible deployment.")

    # 7 Workflow
    slide = new_slide("System Workflow", "End-to-end detection and alert lifecycle")
    steps = [
        ("1", "Packet Input", "Manual scan, structured packet metadata, or tcpdump-style text"),
        ("2", "Validation", "Schema validation and authentication/API key checks"),
        ("3", "Feature Extraction", "Map request fields to flow features and anomaly signals"),
        ("4", "Classification", "Random Forest / detection engine predicts class + confidence"),
        ("5", "Persistence", "Save network log and confidence score in database"),
        ("6", "Alert & Analytics", "Create alert, update dashboard, show recent activity"),
    ]
    for i, (n, title, desc) in enumerate(steps):
        x = 0.75 + i * 2.05
        y = 2.0 if i % 2 == 0 else 3.65
        card(slide, x, y, 1.75, 1.05)
        icon_circle(slide, x + 0.12, y + 0.25, n, BLUE if i < 3 else RED if i == 3 else GREEN, size=0.38)
        add_textbox(slide, x + 0.58, y + 0.18, 1.05, 0.22, title, 9.3, NAVY, True)
        add_textbox(slide, x + 0.18, y + 0.54, 1.42, 0.32, desc, 6.8, SLATE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow(slide, x + 1.75, y + 0.52, x + 2.05, (3.65 if i % 2 == 0 else 2.0) + 0.52, MUTED, 1.3)
    add_textbox(slide, 1.0, 5.45, 10.8, 0.4, "Benign traffic is logged for audit; malicious traffic additionally triggers severity mapping and alert resolution workflow.", 13, NAVY, True, align=PP_ALIGN.CENTER)
    notes.append("Use this slide to explain the viva demo sequence: submit sample traffic, observe prediction, inspect saved log, and review auto-created alert.")

    # 8 Implementation
    slide = new_slide("Implementation", "Important modules, functionalities, and dashboard mockup")
    add_metric(slide, 0.85, 1.45, "Total Traffic", "12,840", BLUE)
    add_metric(slide, 3.55, 1.45, "Threats", "384", RED)
    add_metric(slide, 6.25, 1.45, "Safe Requests", "12,456", GREEN)
    add_metric(slide, 8.95, 1.45, "Active Nodes", "5", CYAN)
    card(slide, 0.85, 2.7, 5.4, 2.95, "Live Traffic Scanner")
    fields = [("Source IP", "192.168.1.50"), ("Destination IP", "10.0.0.1"), ("Protocol", "TCP"), ("Node", "node-1")]
    for i, (label, value) in enumerate(fields):
        x = 1.08 + (i % 2) * 2.45
        y = 3.22 + (i // 2) * 0.72
        add_textbox(slide, x, y, 1.1, 0.18, label.upper(), 6.8, MUTED, True)
        card(slide, x, y + 0.2, 2.0, 0.34, fill=RGBColor(248, 250, 252), line=LINE, radius=False)
        add_textbox(slide, x + 0.08, y + 0.27, 1.8, 0.12, value, 7.5, SLATE, font=MONO_FONT)
    card(slide, 1.08, 4.72, 4.75, 0.5, fill=RGBColor(254, 242, 242), line=RED)
    add_textbox(slide, 1.22, 4.88, 4.3, 0.15, "Threat detected: SYN Flood Attack - 97% confidence", 8.2, RED, True)
    add_code(slide, 6.65, 2.7, 5.7, 2.95, """router.post('/scan', async (req, res) => {
  const parsed = ScanTrafficBody.safeParse(req.body);
  const detection = detectThreat(parsed.data);
  const [log] = await db.insert(networkLogsTable)
    .values({ status: detection.prediction,
              confidenceScore: detection.confidenceScore });
  if (detection.prediction === 'malicious') {
    await db.insert(alertsTable).values({ severity: detection.severity });
  }
});""")
    add_bullets(slide, [
        "Frontend pages: Dashboard, Logs, Alerts, Analytics, Login, Signup.",
        "Backend routes: scan, logs, recent logs, alerts, resolve alert, analytics, summary.",
        "ML modules: dataset generation, model training, feature scaler, label encoder, inference service.",
    ], 0.95, 5.95, 11.4, 0.75, font_size=9.5)
    notes.append("Explain the implemented modules and point out that the dashboard mockup reflects the real React components: summary cards, live scanner, and recent activity.")

    # 9 Results
    slide = new_slide("Results & Testing", "Model evaluation and functional verification")
    add_metric(slide, 0.85, 1.45, "Test Samples", "21,600", BLUE)
    add_metric(slide, 3.55, 1.45, "Accuracy", "100.00%", GREEN)
    add_metric(slide, 6.25, 1.45, "Weighted F1", "1.0000", CYAN)
    add_metric(slide, 8.95, 1.45, "Classes", "7", RED)
    chart_data = CategoryChartData()
    chart_data.categories = ["BENIGN", "DDoS", "GoldenEye", "Hulk", "FTP", "PortScan", "SSH"]
    chart_data.add_series("Support", (14000, 1600, 1000, 2000, 800, 1600, 600))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.95), Inches(2.75), Inches(5.75), Inches(2.75), chart_data).chart
    chart.has_legend = False
    chart.value_axis.visible = True
    chart.category_axis.tick_labels.font.size = Pt(7)
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Test Set Distribution"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    card(slide, 7.25, 2.75, 4.95, 2.75, "Test Cases")
    add_bullets(slide, [
        "Signup/login returns JWT and protects user-specific endpoints.",
        "Scan endpoint validates payload and returns prediction, confidence, threat type, and alert ID.",
        "Logs endpoints support pagination, status filters, search, and recent activity.",
        "Alert endpoint lists severity-filtered alerts and supports resolve action.",
        "ML report shows perfect confusion matrix on synthetic CICIDS-style test data.",
    ], 7.5, 3.18, 4.45, 1.85, font_size=9.2)
    notes.append("Present the results carefully as project/test-dataset results. Clarify that the model was evaluated on the generated CICIDS-style dataset included with the project.")

    # 10 Advantages
    slide = new_slide("Advantages & Applications", "Benefits and real-world use cases")
    left_items = [
        ("Real-time Awareness", "Auto-refreshing dashboard and recent logs."),
        ("Actionable Alerts", "Malicious traffic becomes a severity-ranked incident."),
        ("Auditability", "Every scan stores source, destination, status, confidence, and node."),
        ("Scalable Design", "Node-based input and separate UI/API/database layers."),
    ]
    right_items = [
        ("Campus Networks", "Monitor lab, department, and server traffic."),
        ("Small SOC Teams", "Central dashboard for security operators."),
        ("Cloud Workloads", "Detect suspicious API or service traffic patterns."),
        ("Learning Platform", "Demonstrate IDS concepts, ML classification, and incident workflow."),
    ]
    add_textbox(slide, 1.0, 1.45, 4.3, 0.28, "Advantages", 15, NAVY, True)
    for i, (t, d) in enumerate(left_items):
        y = 1.95 + i * 0.88
        icon_circle(slide, 1.0, y, "+", GREEN, size=0.34)
        add_textbox(slide, 1.45, y - 0.02, 3.5, 0.22, t, 11, NAVY, True)
        add_textbox(slide, 1.45, y + 0.25, 4.4, 0.2, d, 8.5, SLATE)
    add_textbox(slide, 7.0, 1.45, 4.3, 0.28, "Applications", 15, NAVY, True)
    for i, (t, d) in enumerate(right_items):
        y = 1.95 + i * 0.88
        icon_circle(slide, 7.0, y, ">", BLUE, size=0.34)
        add_textbox(slide, 7.45, y - 0.02, 3.5, 0.22, t, 11, NAVY, True)
        add_textbox(slide, 7.45, y + 0.25, 4.4, 0.2, d, 8.5, SLATE)
    card(slide, 1.0, 5.65, 10.9, 0.62, fill=RGBColor(239, 246, 255), line=BLUE)
    add_textbox(slide, 1.3, 5.84, 10.2, 0.2, "The system converts raw network activity into classified evidence that can be reviewed, explained, and acted upon.", 11, BLUE, True, align=PP_ALIGN.CENTER)
    notes.append("Connect benefits to practical environments. For viva, emphasize how each benefit is visible in the implemented application.")

    # 11 Future
    slide = new_slide("Future Enhancements", "Improvements and scalability options")
    future = [
        ("Deeper Packet Capture", "Integrate Zeek or Suricata agents for continuous real traffic streaming."),
        ("Distributed Agents", "Deploy lightweight node agents that stream features to the central API."),
        ("Model Maturity", "Train on larger real datasets, tune thresholds, and measure false positives."),
        ("Response Automation", "Add blocking rules, ticket creation, email/SMS alerts, and playbooks."),
        ("Role-Based Access", "Separate admin, analyst, and viewer permissions."),
        ("Production Observability", "Add metrics, tracing, model drift monitoring, and backup strategy."),
    ]
    for i, (t, d) in enumerate(future):
        x = 0.85 + (i % 3) * 4.1
        y = 1.55 + (i // 3) * 2.0
        card(slide, x, y, 3.45, 1.45)
        icon_circle(slide, x + 0.18, y + 0.25, str(i + 1), [BLUE, CYAN, GREEN, RED, AMBER, SLATE][i], size=0.42)
        add_textbox(slide, x + 0.75, y + 0.22, 2.35, 0.25, t, 10.8, NAVY, True)
        add_textbox(slide, x + 0.3, y + 0.72, 2.85, 0.45, d, 8.3, SLATE, align=PP_ALIGN.CENTER)
    notes.append("Acknowledge that the academic prototype is functional, then show a credible path toward a production-grade DIDS.")

    # 12 Conclusion
    slide = new_slide("Conclusion", "Summary and key achievements")
    card(slide, 0.95, 1.55, 5.45, 3.65, "Summary")
    add_bullets(slide, [
        "Developed a full-stack distributed intrusion detection dashboard.",
        "Implemented authentication, scan submission, logs, alerts, analytics, and summary APIs.",
        "Integrated trained Random Forest artifacts and detection logic for threat classification.",
        "Stored results persistently with confidence scores and alert metadata.",
    ], 1.2, 2.05, 4.85, 2.55, font_size=12.3)
    card(slide, 7.0, 1.55, 4.95, 3.65, "Key Achievements")
    add_bullets(slide, [
        "Seven-class threat coverage: BENIGN, DDoS, DoS Hulk, DoS GoldenEye, PortScan, FTP-Patator, SSH-Patator.",
        "100% accuracy and weighted F1 on the included synthetic CICIDS-style test set.",
        "University defense-ready workflow with reproducible Docker-based deployment.",
    ], 7.25, 2.05, 4.35, 2.35, font_size=12.3)
    add_textbox(slide, 1.1, 5.85, 10.9, 0.35, "Final takeaway: the project demonstrates how ML-assisted detection can be combined with usable dashboards and reliable backend services.", 13, BLUE, True, align=PP_ALIGN.CENTER)
    notes.append("Close the technical story. Restate the main achievement: a complete system, not only a standalone model.")

    # 13 Thank you
    slide = new_slide()
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    icon_circle(slide, 6.2, 1.28, "D", BLUE, size=0.9)
    add_textbox(slide, 3.25, 2.35, 6.9, 0.75, "Thank You", 42, WHITE, True, align=PP_ALIGN.CENTER, font=TITLE_FONT)
    add_textbox(slide, 3.25, 3.15, 6.9, 0.4, "Questions & Answers", 20, RGBColor(191, 219, 254), True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 3.25, 4.25, 6.9, 0.35, "Muhammad Zaid (232674)  |  M. Shafiq (232699)  |  GCUF", 13, RGBColor(226, 232, 240), align=PP_ALIGN.CENTER)
    notes.append("Invite questions. Be prepared to discuss architecture, detection logic, test results, limitations, and future improvements.")

    prs.save(OUT)
    add_slide_notes(OUT, notes)
    return OUT


if __name__ == "__main__":
    path = build_deck()
    print(path)
